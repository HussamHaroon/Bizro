# Bizro pitch deck generator — "stamped-ledger neobrutalism".
# Facts sourced from site/src/content.ts only. Run: python presentation/build_pitch.py
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

CREAM = RGBColor(0xF5, 0xF1, 0xE6)
INK = RGBColor(0x1F, 0x1B, 0x16)
GREEN = RGBColor(0x0B, 0x5D, 0x3B)
RED = RGBColor(0xA6, 0x33, 0x2B)
GOLD = RGBColor(0xE9, 0xA9, 0x3D)
TEAL = RGBColor(0x1F, 0x7A, 0x6C)
PAPER = RGBColor(0xFB, 0xF8, 0xF0)  # frame interior, one step lighter than bg

SLAB = "Arial Black"
BODY = "Arial"

SHADOW_OFF = Inches(0.08)
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _no_autoshadow(shape):
    shape.shadow.inherit = False


def frame(s, x, y, w, h, fill=PAPER, line=INK, line_w=3.0, shape=MSO_SHAPE.RECTANGLE):
    """Hard offset shadow (ink duplicate, +0.08in, zero blur) then the framed block."""
    sh = s.shapes.add_shape(shape, int(x + SHADOW_OFF), int(y + SHADOW_OFF), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = INK
    sh.line.fill.background()
    _no_autoshadow(sh)
    f = s.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    f.fill.solid()
    f.fill.fore_color.rgb = fill
    f.line.color.rgb = line
    f.line.width = Pt(line_w)
    _no_autoshadow(f)
    return f


def bar(s, x, y, w, h, color):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.color.rgb = INK
    b.line.width = Pt(1.5)
    _no_autoshadow(b)
    return b


def text(s, x, y, w, h, runs, size=18, color=INK, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.0, wrap=True):
    """runs: str, or list of paragraphs; each paragraph is str or list of
    (txt, {size,color,font,bold}) run tuples."""
    tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if isinstance(para, str):
            para = [(para, {})]
        for txt, ov in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(ov.get("size", size))
            r.font.color.rgb = ov.get("color", color)
            r.font.name = ov.get("font", font)
            r.font.bold = ov.get("bold", bold)
    return tb


def chip(s, x, y, w, label, color, txt_color=None, size=12):
    c = bar(s, x, y, w, Inches(0.34), color)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.name = BODY
    r.font.bold = True
    r.font.color.rgb = txt_color if txt_color else CREAM
    return c


def headline(s, title, sub=None, chip_label=None, chip_color=GREEN,
             gold_underline=False, y=Inches(0.42)):
    """Standard slide head: optional chip, slab headline, gold underline bar."""
    x = Inches(0.6)
    if chip_label:
        chip(s, x, y, Inches(0.28 + 0.09 * len(chip_label)), chip_label, chip_color)
        y = int(y + Inches(0.52))
    size = 30 if len(title) > 60 else 34
    cpl = max(1, 11.95 * 72 / (0.66 * size))  # Arial Black avg char width
    lines = max(1, -(-len(title) // int(cpl)))
    line_h = size * 1.30 / 72
    text(s, x, y, Inches(12.1), Inches(lines * line_h + 0.12), title,
         size=size, font=SLAB, leading=0.98)
    uy = int(y + Inches(line_h * lines + 0.10))
    bar(s, x, uy, Inches(2.2), Inches(0.09), GOLD if gold_underline else INK)
    if sub:
        text(s, x, int(uy + Inches(0.16)), Inches(12.1), Inches(0.4), sub, size=15)
    return int(uy + Inches(0.55 if sub else 0.35))


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ---------------------------------------------------------------- 1 · TITLE
s = new_slide()
frame(s, Inches(1.2), Inches(1.55), Inches(10.9), Inches(4.1), fill=PAPER)
# gold seal circle, top-right
seal = frame(s, Inches(11.3), Inches(0.45), Inches(1.45), Inches(1.45),
             fill=GOLD, line=INK, line_w=3.0, shape=MSO_SHAPE.OVAL)
tf = seal.text_frame
tf.word_wrap = False
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "BIZRO"
r.font.size = Pt(13); r.font.name = SLAB; r.font.color.rgb = INK
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "2026"
r.font.size = Pt(11); r.font.name = BODY; r.font.bold = True; r.font.color.rgb = INK

text(s, Inches(1.5), Inches(1.85), Inches(10.3), Inches(1.9), "Bizro",
     size=110, font=SLAB, align=PP_ALIGN.CENTER)
bar(s, Inches(5.55), Inches(3.85), Inches(2.2), Inches(0.12), INK)
text(s, Inches(1.5), Inches(4.15), Inches(10.3), Inches(0.6),
     "The paper khata, given a memory.", size=26, align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.5),
     "A zero-typing Voice & Vision copilot for Pakistan's micro-entrepreneurs",
     size=14, align=PP_ALIGN.CENTER, color=TEAL, bold=True)
frame(s, Inches(1.2), Inches(6.35), Inches(10.9), Inches(0.7), fill=INK, line=INK)
text(s, Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.5),
     "Bano Qabil × Alibaba Cloud AI Hackathon Pakistan 2026",
     size=15, color=CREAM, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Open on the name and the promise: Bizro gives the paper khata a memory. "
         "One sentence on who it is for — the Pakistani micro-entrepreneur who already "
         "runs his business on WhatsApp. Then move straight into the problem.")

# -------------------------------------------------------------- 2 · PROBLEM
s = new_slide()
top = headline(s, "Creditworthy, but invisible to the lender meant for them.",
               chip_label="01 · The problem")
stats = [
    ("Formal account", RED, "10.3%", 54,
     "of Pakistani adults hold a formal financial-institution account — at the last national baseline."),
    ("South Asia, for contrast", TEAL, "~33%", 54,
     "the South Asian average of adults with a formal financial-institution account."),
    ("The blocker", GOLD, "Shariah-\ncompliant\ndemand", 30,
     "small businesses avoiding formal finance cite one specific reason: they are waiting for a Shariah-compliant option, not an interest-bearing one."),
]
fx, fw, gap = Inches(0.6), Inches(3.85), Inches(0.24)
for i, (cchip, tone, num, nsize, body) in enumerate(stats):
    x = int(fx + i * (fw + gap))
    frame(s, x, top, fw, Inches(4.35))
    chip(s, int(x + Inches(0.25)), int(top + Inches(0.25)),
         Inches(0.3 + 0.09 * len(cchip)), cchip, tone,
         txt_color=INK if tone == GOLD else CREAM)
    text(s, int(x + Inches(0.25)), int(top + Inches(0.85)), int(fw - Inches(0.5)),
         Inches(1.9), num.replace("\n", " ") if nsize < 40 else num,
         size=nsize, font=SLAB, color=tone, leading=0.95)
    text(s, int(x + Inches(0.25)), int(top + Inches(2.55)), int(fw - Inches(0.5)),
         Inches(1.9), body, size=13.5, leading=1.05)
text(s, Inches(0.6), int(top + Inches(4.5)), Inches(12.1), Inches(0.35),
     [[("The gap isn't creditworthiness — it's legibility.", {"bold": True, "size": 16})]],
     align=PP_ALIGN.CENTER)
notes(s, "Only about one in ten Pakistani adults has a formal financial account, against a "
         "South Asian average near a third. And the blocker isn't money — a meaningful share of "
         "small businesses say they are waiting for a Shariah-compliant option. These people are "
         "creditworthy; they are just invisible on paper.")

# -------------------------------------------------------- 3 · THE RAIL EXISTS
s = new_slide()
top = headline(s, "The interest-free rail already exists",
               chip_label="Mawakhat · Alkhidmat Foundation", gold_underline=True)
frame(s, Inches(0.6), top, Inches(12.13), Inches(1.5))
text(s, Inches(0.9), int(top + Inches(0.2)), Inches(11.5), Inches(1.2),
     "Alkhidmat's Mawakhat program is Qarz-e-Hasna (interest-free) microfinance funded "
     "through zakat and sadaqat — the credit product these shopkeepers would actually accept.",
     size=17, leading=1.15)
minis = [("~800", "branches across 400+ cities", GREEN),
         ("PKR 30–75k", "typical Qarz-e-Hasna loan", TEAL),
         ("99.9%", "repayment rate*", RED)]
mw = Inches(3.85)
for i, (dt, dd, tone) in enumerate(minis):
    x = int(Inches(0.6) + i * (mw + Inches(0.24)))
    y = int(top + Inches(1.8))
    frame(s, x, y, mw, Inches(1.7))
    bar(s, x, y, Inches(0.14), Inches(1.7), tone)
    text(s, int(x + Inches(0.35)), int(y + Inches(0.2)), int(mw - Inches(0.6)),
         Inches(0.7), dt, size=30, font=SLAB, color=tone)
    text(s, int(x + Inches(0.35)), int(y + Inches(0.95)), int(mw - Inches(0.6)),
         Inches(0.6), dd, size=14)
frame(s, Inches(0.6), int(top + Inches(3.8)), Inches(12.13), Inches(1.15), fill=INK, line=INK)
text(s, Inches(0.9), int(top + Inches(3.98)), Inches(11.5), Inches(0.85),
     "What it can't see is a business history that lives in a handwritten notebook.",
     size=20, font=SLAB, color=CREAM, align=PP_ALIGN.CENTER, leading=1.0)
text(s, Inches(0.6), int(top + Inches(5.1)), Inches(12.1), Inches(0.3),
     "*as claimed by Mawakhat", size=11)
notes(s, "The lending rail already exists: Alkhidmat's Mawakhat program runs Qarz-e-Hasna, "
         "interest-free microfinance funded by zakat — roughly 800 branches in 400-plus cities, "
         "loans of 30 to 75 thousand rupees, and a claimed 99.9 percent repayment rate. The one "
         "thing Mawakhat cannot see is the shopkeeper's business history, because it lives in a "
         "handwritten notebook.")

# ------------------------------------------------------------- 4 · SOLUTION
s = new_slide()
top = headline(s, "Bizro, in one line", chip_label="The solution")
frame(s, Inches(0.6), top, Inches(12.13), Inches(1.55), fill=GREEN, line=INK)
text(s, Inches(0.95), int(top + Inches(0.18)), Inches(11.4), Inches(1.25),
     "A zero-typing Voice & Vision copilot that turns WhatsApp voice notes and "
     "receipt photos into a lender-legible credit history.",
     size=21, font=SLAB, color=CREAM, leading=1.05)
steps = [
    ("VOICE", GREEN, "Send the voice note", "Parsed into a structured sale or credit entry; a stamped invoice returns on WhatsApp.", "Qwen3.5-Omni-Plus"),
    ("VISION", TEAL, "Photograph the receipt", "The expense is logged — and obvious pricing errors are flagged.", "Qwen-VL-OCR"),
    ("REPORT", GOLD, "Tap, months later", "Accumulated history becomes a Mawakhat-style Credit Readiness Report.", "Qwen3.7-Plus"),
]
sw_, sy = Inches(3.7), int(top + Inches(1.95))
for i, (name, tone, h3, body, model) in enumerate(steps):
    x = int(Inches(0.6) + i * (sw_ + Inches(0.55)))
    frame(s, x, sy, sw_, Inches(2.9))
    chip(s, int(x + Inches(0.22)), int(sy + Inches(0.22)), Inches(1.35), name, tone,
         txt_color=INK if tone == GOLD else CREAM, size=13)
    text(s, int(x + Inches(0.22)), int(sy + Inches(0.72)), int(sw_ - Inches(0.44)),
         Inches(0.55), h3, size=17, font=SLAB)
    text(s, int(x + Inches(0.22)), int(sy + Inches(1.3)), int(sw_ - Inches(0.44)),
         Inches(1.1), body, size=13, leading=1.08)
    text(s, int(x + Inches(0.22)), int(sy + Inches(2.45)), int(sw_ - Inches(0.44)),
         Inches(0.35), model, size=11.5, bold=True, color=tone if tone != GOLD else RED)
    if i < 2:
        text(s, int(x + sw_ + Inches(0.06)), int(sy + Inches(1.1)), Inches(0.5),
             Inches(0.6), "→", size=30, font=SLAB, align=PP_ALIGN.CENTER)
text(s, Inches(0.6), int(sy + Inches(3.1)), Inches(12.1), Inches(0.4),
     "No typing. No new app. No new habit.", size=15, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Bizro is a zero-typing Voice and Vision copilot. He sends the WhatsApp voice note he "
         "was sending anyway; he photographs the receipt he already keeps. Bizro turns both into "
         "structured ledger entries — and months of those become a credit history a lender like "
         "Mawakhat can actually read.")

# ------------------------------------------------------------ 5 · LIVE DEMO
s = new_slide()
top = headline(s, "This is running live right now", chip_label="05 · Live demo",
               chip_color=RED)
frame(s, Inches(0.6), top, Inches(7.3), Inches(2.5), fill=INK, line=INK)
text(s, Inches(0.95), int(top + Inches(0.3)), Inches(6.6), Inches(0.6),
     "https://bizro-pk.vercel.app", size=24, font=SLAB, color=GOLD)
text(s, Inches(0.95), int(top + Inches(1.1)), Inches(6.6), Inches(1.2),
     ["Record a voice note.", "Watch it become a ledger entry."],
     size=19, color=CREAM, leading=1.25)
frame(s, Inches(8.2), top, Inches(4.5), Inches(2.5), fill=PAPER)
text(s, Inches(8.4), int(top + Inches(0.95)), Inches(4.1), Inches(0.7),
     "[ screenshot space ]", size=14, align=PP_ALIGN.CENTER, color=TEAL, bold=True)
frame(s, Inches(0.6), int(top + Inches(2.85)), Inches(12.1), Inches(2.35), fill=PAPER)
text(s, Inches(0.95), int(top + Inches(3.05)), Inches(11.4), Inches(0.4),
     "DEMO PATH", size=13, font=SLAB, color=RED)
text(s, Inches(0.95), int(top + Inches(3.5)), Inches(11.4), Inches(1.6),
     ["1 · Open the live simulator — a WhatsApp-style chat with Bizro.",
      "2 · Record an Urdu voice note, or drop in a receipt photo.",
      "3 · Watch the pipeline parse it, stamp it, and log it to the ledger."],
     size=15, leading=1.35)
notes(s, "This is not a mock-up — the product is deployed and running live at this URL right now. "
         "In the simulator you record an Urdu voice note or upload a receipt photo, and the real "
         "pipeline parses it, stamps it, and writes it into the ledger. If time allows we will do "
         "this on stage.")

# --------------------------------------------------------- 6 · HOW IT WORKS
s = new_slide()
top = headline(s, "Three taps' worth of effort — months' worth of credit history.",
               chip_label="02 · How it works")
steps = [
    ("Step 1 · Voice", GREEN, "Send the voice note",
     "A casual Urdu voice note, the way he'd talk to an employee. It's parsed into a structured sale or credit entry, and a stamped invoice comes straight back on WhatsApp.",
     "Qwen3.5-Omni-Plus · WhatsApp in / out"),
    ("Step 2 · Vision", TEAL, "Photograph the receipt",
     "One photo of the messy, handwritten supplier receipt. The expense is logged — and obvious pricing errors are flagged before they can hide in a hand-tallied notebook.",
     "Qwen-VL-OCR · Price-error flags"),
    ("Step 3 · Report", GOLD, "Tap, months later",
     "One tap turns months of accumulated, source-linked history into a Mawakhat-style Credit Readiness Report — with a full audit trail a loan officer can drill into.",
     "Qwen3.7-Plus · Mawakhat-style format"),
]
cw, cy = Inches(3.95), int(top + Inches(0.1))
for i, (cchip, tone, h3, body, badge) in enumerate(steps):
    x = int(Inches(0.6) + i * (cw + Inches(0.14)))
    frame(s, x, cy, cw, Inches(4.3))
    chip(s, int(x + Inches(0.22)), int(cy + Inches(0.22)),
         Inches(0.3 + 0.09 * len(cchip)), cchip, tone,
         txt_color=INK if tone == GOLD else CREAM)
    text(s, int(x + Inches(0.22)), int(cy + Inches(0.75)), int(cw - Inches(0.44)),
         Inches(0.85), h3, size=19, font=SLAB, leading=0.95)
    text(s, int(x + Inches(0.22)), int(cy + Inches(1.6)), int(cw - Inches(0.44)),
         Inches(1.95), body, size=13.5, leading=1.15)
    bar(s, int(x + Inches(0.22)), int(cy + Inches(3.6)), int(cw - Inches(0.44)),
        Inches(0.03), INK)
    text(s, int(x + Inches(0.22)), int(cy + Inches(3.7)), int(cw - Inches(0.44)),
         Inches(0.5), badge, size=11, bold=True, color=tone if tone != GOLD else RED)
text(s, Inches(0.6), int(cy + Inches(4.5)), Inches(12.1), Inches(0.4),
     "Zero typing, all of it.", size=15, bold=True, align=PP_ALIGN.CENTER)
notes(s, "The whole flow is three taps' worth of effort. Step one: he sends a voice note and gets "
         "a stamped invoice back. Step two: one photo of a handwritten receipt logs the expense "
         "and flags pricing errors. Step three, months later: one tap produces the Credit "
         "Readiness Report in a format Mawakhat loan officers already understand.")

# ------------------------------------------------------------ 7 · TECHNOLOGY
s = new_slide()
top = headline(s, "Provider-agnostic AI backbone", chip_label="Technology")
pipe = [
    ("1 · Speech-to-text", GREEN, "Urdu / mixed-language voice note → transcript. Runs on Qwen3.5-Omni-Plus."),
    ("2 · Structured JSON", TEAL, "Transcript & receipt photo → typed transaction JSON with confidence scores. Qwen-VL-OCR reads the receipts."),
    ("3 · Urdu narrative", GOLD, "Ledger history → a Credit Readiness Report a loan officer can read, reasoned by Qwen3.7-Plus."),
]
cw, cy = Inches(3.7), int(top + Inches(0.1))
for i, (h3, tone, body) in enumerate(pipe):
    x = int(Inches(0.6) + i * (cw + Inches(0.55)))
    frame(s, x, cy, cw, Inches(2.6))
    bar(s, x, cy, cw, Inches(0.16), tone)
    text(s, int(x + Inches(0.22)), int(cy + Inches(0.35)), int(cw - Inches(0.44)),
         Inches(0.6), h3, size=17, font=SLAB, leading=0.95)
    text(s, int(x + Inches(0.22)), int(cy + Inches(1.05)), int(cw - Inches(0.44)),
         Inches(1.45), body, size=13, leading=1.12)
    if i < 2:
        text(s, int(x + cw + Inches(0.06)), int(cy + Inches(0.95)), Inches(0.5),
             Inches(0.6), "→", size=30, font=SLAB, align=PP_ALIGN.CENTER)
frame(s, Inches(0.6), int(cy + Inches(3.0)), Inches(12.1), Inches(1.35), fill=INK, line=INK)
text(s, Inches(0.95), int(cy + Inches(3.2)), Inches(11.4), Inches(1.0),
     [[("Runs on free tiers today. ", {"color": CREAM, "bold": True}),
       ("Qwen via Alibaba Cloud Model Studio is the production path — ", {"color": CREAM}),
       ("one environment variable.", {"color": CREAM, "bold": True, "font": SLAB})]],
     size=18, leading=1.2)
text(s, Inches(0.6), int(cy + Inches(4.6)), Inches(12.1), Inches(0.6),
     "Stack: FastAPI + Postgres-compatible storage · React dashboard · WhatsApp Cloud API shape · deployed on Vercel.",
     size=13, align=PP_ALIGN.CENTER)
notes(s, "The backbone is deliberately provider-agnostic. Voice goes to speech-to-text, then to "
         "structured JSON with confidence scores, then to an Urdu narrative for loan officers. "
         "Today it runs entirely on free tiers; switching to Qwen through Alibaba Cloud Model "
         "Studio — the production path — is literally one environment variable.")

# ---------------------------------------------------------- 8 · TRUST & AUDIT
s = new_slide()
top = headline(s, "Every number traces to a voice note or a photo.",
               chip_label="04 · Trust & audit", chip_color=TEAL)
frame(s, Inches(0.6), top, Inches(6.0), Inches(4.5))
text(s, Inches(0.9), int(top + Inches(0.25)), Inches(5.4), Inches(0.4),
     "THE AUDIT TRAIL IS THE PRODUCT", size=13, font=SLAB, color=RED)
text(s, Inches(0.9), int(top + Inches(0.75)), Inches(5.4), Inches(1.6),
     "In the live app, every field drills down to the original voice note or receipt photo behind it.",
     size=16, leading=1.2)
for j, (h3, body) in enumerate([
        ("Source + confidence, always", "Every AI-parsed line keeps its source and a confidence score. A visible trail, not a black box."),
        ("One-tap correct", "Any entry can be fixed in one tap — and the correction is itself a trust signal."),
]):
    yy = int(top + Inches(2.15) + j * Inches(1.2))
    bar(s, Inches(0.9), yy, Inches(0.1), Inches(1.0), GREEN if j == 0 else TEAL)
    text(s, Inches(1.15), yy, Inches(5.15), Inches(0.35), h3, size=14, font=SLAB)
    text(s, Inches(1.15), int(yy + Inches(0.38)), Inches(5.15), Inches(0.7), body,
         size=12.5, leading=1.08)
frame(s, Inches(6.95), top, Inches(5.75), Inches(4.5), fill=PAPER)
text(s, Inches(7.25), int(top + Inches(0.22)), Inches(5.15), Inches(0.4),
     "Ledger entry · 12 Aug", size=13, font=SLAB)
chip(s, Inches(11.3), int(top + Inches(0.2)), Inches(1.1), "UDHAR · ادھار", RED, size=11)
text(s, Inches(7.25), int(top + Inches(0.75)), Inches(5.15), Inches(0.5),
     "Ahmad — credit given", size=19, font=SLAB)
text(s, Inches(7.25), int(top + Inches(1.35)), Inches(5.15), Inches(0.55),
     [[("PKR 5,000  ", {"size": 26, "font": SLAB, "color": GREEN}),
       ("پانچ ہزار", {"size": 16, "color": TEAL})]])
for k, line in enumerate([
        "Source: WhatsApp voice note · 0:14",
        "Parsed by Qwen3.5-Omni-Plus · confidence 96%",
        "One-tap correct — the fix is itself a trust signal"]):
    yy = int(top + Inches(2.15) + k * Inches(0.52))
    bar(s, Inches(7.25), int(yy + Inches(0.4)), Inches(5.1), Inches(0.02), INK)
    text(s, Inches(7.25), yy, Inches(5.15), Inches(0.4), line, size=13)
st = frame(s, Inches(9.6), int(top + Inches(3.75)), Inches(2.9), Inches(0.6),
           fill=GOLD, line=INK)
text(s, Inches(9.6), int(top + Inches(3.85)), Inches(2.9), Inches(0.4),
     "AI-Parsed · Confirmed", size=13, font=SLAB, align=PP_ALIGN.CENTER)
notes(s, "Trust is the product, not a feature. Every parsed entry keeps its source — the actual "
         "voice note or receipt photo — plus a confidence score, and any mistake is one tap from "
         "corrected. That visible trail is what lets a loan officer trust the report as much as "
         "the shopkeeper trusts the tool.")

# ----------------------------------------------------------- 9 · WHY BIZRO
s = new_slide()
top = headline(s, "More than \u201cOCR plus a chatbot.\u201d", chip_label="03 · Why Bizro")
cards = [
    ("Udhar Radar", GREEN, "Money owed TO the shop",
     "Flips the expense-tracker lens: Bizro tracks the money customers owe to the shopkeeper — the dominant real use of a paper khata, and the piece most digitization tools miss entirely."),
    ("Audit trail on every entry", TEAL, "Source + confidence, always",
     "Every AI-parsed line keeps its source voice note or photo plus a confidence score — and a one-tap correct. A visible trail, not a black box."),
    ("A direct line to Mawakhat", GOLD, "A real credit rail, not a score",
     "Not a generic \u201ccredit score.\u201d The report is shaped for a lending program that already exists, already has ~800 branches, and already has an institutional reason to want exactly this evidence."),
]
cw, cy = Inches(3.95), int(top + Inches(0.1))
for i, (h3, tone, tag, body) in enumerate(cards):
    x = int(Inches(0.6) + i * (cw + Inches(0.14)))
    frame(s, x, cy, cw, Inches(4.7))
    bar(s, x, cy, Inches(0.14), Inches(4.7), tone if tone != GOLD else INK)
    text(s, int(x + Inches(0.35)), int(cy + Inches(0.25)), int(cw - Inches(0.6)),
         Inches(0.95), h3, size=18, font=SLAB, leading=0.95)
    chip(s, int(x + Inches(0.35)), int(cy + Inches(1.2)),
         Inches(0.3 + 0.085 * len(tag)), tag, tone,
         txt_color=INK if tone == GOLD else CREAM, size=11)
    text(s, int(x + Inches(0.35)), int(cy + Inches(1.75)), int(cw - Inches(0.6)),
         Inches(2.8), body, size=13.5, leading=1.15)
text(s, Inches(0.6), int(cy + Inches(4.9)), Inches(12.1), Inches(0.4),
     "Built for how karyana shops actually run.", size=15, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Three things separate Bizro from a generic OCR-plus-chatbot. Udhar Radar tracks money "
         "owed TO the shop — the real job of a paper khata, which digitization tools miss. Every "
         "entry carries its audit trail. And the report is shaped for Mawakhat specifically: a "
         "lender that already exists and already wants this evidence.")

# -------------------------------------------------------------- 10 · IMPACT
s = new_slide()
top = headline(s, "Months of voice notes become a credit history a lender can actually read.",
               chip_label="Impact")
frame(s, Inches(0.6), top, Inches(12.13), Inches(1.9), fill=GREEN, line=INK)
text(s, Inches(0.95), int(top + Inches(0.22)), Inches(11.4), Inches(1.6),
     [[("For the ", {"color": CREAM, "size": 24}),
       ("10 million+", {"color": GOLD, "size": 40, "font": SLAB}),
       (" unbanked micro-entrepreneurs of Pakistan.", {"color": CREAM, "size": 24})]],
     leading=1.1)
caps = [
    "He sends the voice note. He snaps the receipt. That's the whole job.",
    "No typing. The pieces assemble themselves into a ledger entry.",
    "Every entry is stamped — parsed, priced, and double-checked.",
]
for j, cline in enumerate(caps):
    yy = int(top + Inches(2.1) + j * Inches(0.78))
    frame(s, Inches(0.6), yy, Inches(12.13), Inches(0.7))
    text(s, Inches(0.95), int(yy + Inches(0.13)), Inches(11.4), Inches(0.45),
         cline, size=15, bold=True)
text(s, Inches(0.6), int(top + Inches(4.5)), Inches(12.1), Inches(0.4),
     "The behavior already exists. Bizro just makes it legible.",
     size=17, font=SLAB, align=PP_ALIGN.CENTER, color=RED)
notes(s, "The impact: months of voice notes he was already sending become a credit history a "
         "lender can read — for the ten million plus unbanked micro-entrepreneurs of Pakistan. "
         "No behavior change is required; the voice note and the receipt photo already exist. "
         "Bizro simply makes them legible to the rail that was built for him.")

# -------------------------------------------------------- 11 · THANK YOU
s = new_slide()
frame(s, Inches(1.2), Inches(1.3), Inches(10.9), Inches(4.4), fill=PAPER)
seal = frame(s, Inches(11.3), Inches(0.45), Inches(1.45), Inches(1.45),
             fill=GOLD, line=INK, line_w=3.0, shape=MSO_SHAPE.OVAL)
tf = seal.text_frame; tf.word_wrap = False
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "شکریہ"
r.font.size = Pt(15); r.font.name = BODY; r.font.bold = True; r.font.color.rgb = INK

text(s, Inches(1.5), Inches(1.75), Inches(10.3), Inches(1.4), "Thank you.",
     size=76, font=SLAB, align=PP_ALIGN.CENTER)
bar(s, Inches(5.55), Inches(3.35), Inches(2.2), Inches(0.12), INK)
text(s, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.5),
     "Team Bizro · Bano Qabil × Alibaba Cloud AI Hackathon Pakistan 2026",
     size=16, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
     "https://bizro-pk.vercel.app", size=24, font=SLAB, color=GREEN,
     align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.4),
     "Demo available on request.", size=15, align=PP_ALIGN.CENTER, color=TEAL, bold=True)
frame(s, Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.8), fill=INK, line=INK)
text(s, Inches(1.5), Inches(6.38), Inches(10.3), Inches(0.5),
     "The paper khata, given a memory.", size=16, color=CREAM, bold=True,
     align=PP_ALIGN.CENTER)
notes(s, "Close on the tagline: the paper khata, given a memory. The product is live at this URL "
         "today, and a full demo is available on request. Thank the judges and invite questions.")

OUT = "presentation/Bizro_Pitch.pptx"
prs.save(OUT)
print("saved", OUT)
