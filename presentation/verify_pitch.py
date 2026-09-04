import math, os
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
prs = Presentation("presentation/Bizro_Pitch.pptx")
SW = prs.slide_width / EMU_IN
SH = prs.slide_height / EMU_IN
print(f"slides: {len(prs.slides)}  size: {SW:.3f} x {SH:.3f} in")
issues = 0
for idx, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        l, t = sh.left / EMU_IN, sh.top / EMU_IN
        w, h = sh.width / EMU_IN, sh.height / EMU_IN
        if l < -0.01 or t < -0.01 or l + w > SW + 0.01 or t + h > SH + 0.01:
            print(f"  [S{idx}] BOUNDS {sh.shape_type} '{getattr(sh,'name','')}' "
                  f"l={l:.2f} t={t:.2f} r={l+w:.2f} b={t+h:.2f}")
            issues += 1
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        est_h = 0.0
        for p in tf.paragraphs:
            runs = [(r.text, (r.font.size.pt if r.font.size else 18),
                     r.font.name or "Arial") for r in p.runs if r.text]
            if not runs:
                continue
            avail = max(0.5, w - 0.15)
            max_size = max(sz for _, sz, _ in runs)
            # split runs on explicit newlines into visual paragraphs
            groups, cur = [], []
            for txt, sz, fn in runs:
                parts = txt.split("\n")
                for i, part in enumerate(parts):
                    if i > 0:
                        groups.append(cur); cur = []
                    if part:
                        cur.append((part, sz, fn))
            groups.append(cur)
            ls = p.line_spacing or 1.0
            for g in groups:
                if not g:
                    est_h += max_size * 1.22 * ls / 72
                    continue
                total_w = sum(len(t) * (0.66 if f == "Arial Black" else 0.52) * sz / 72
                              for t, sz, f in g)
                gmax = max(sz for _, sz, _ in g)
                nlines = max(1, math.ceil(total_w / avail))
                est_h += nlines * gmax * 1.22 * ls / 72
        if est_h > h + 0.12 and tf.word_wrap is not False:
            print(f"  [S{idx}] TEXT-OVERFLOW? box h={h:.2f}in est={est_h:.2f}in :: "
                  f"'{tf.text[:60]}'")
            issues += 1
size_bytes = os.path.getsize("presentation/Bizro_Pitch.pptx")
print(f"file size: {size_bytes} bytes ({size_bytes/1024:.1f} KB)")
print("OK — no issues" if issues == 0 else f"{issues} potential issue(s) above")
